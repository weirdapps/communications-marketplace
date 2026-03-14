#!/usr/bin/env python3
"""
NBG Presentation Builder

Creates NBG-branded presentations from scratch using python-pptx.
Every slide is built programmatically with exact NBG brand specifications:
- White backgrounds on ALL slides
- Aptos font throughout, margin: 0, valign: top
- Pixel-perfect positioning per the brand system
- Cyan bullets, proper typography hierarchy
- Logo placement, page numbers on content slides only

Usage:
    python nbg_build.py storyline.yaml output.pptx
"""

import subprocess
import sys
from pathlib import Path

import yaml

# python-pptx and lxml are required for building presentations.
# Imported conditionally so tests that only exercise normalize_slide_type()
# and load_catalog() can run without these heavy dependencies.
try:
    from lxml.etree import SubElement
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt
    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
SCRIPT_DIR = Path(__file__).parent
ASSETS_DIR = SCRIPT_DIR.parent.parent / 'assets'
CATALOG_PATH = ASSETS_DIR / 'slide-catalog.yaml'
TEMPLATES_DIR = ASSETS_DIR / 'templates'

# Logo assets — PNG for python-pptx compatibility (SVG not supported)
LOGO_PNG = ASSETS_DIR / 'nbg-logo-gr.png'
BACK_COVER_LOGO_PNG = ASSETS_DIR / 'nbg-back-cover-logo.png'

# ---------------------------------------------------------------------------
# NBG Brand Constants
# ---------------------------------------------------------------------------
FONT = 'Aptos'
FONT_BULLET = 'Arial'

if _HAS_PPTX:
    # Colors (RGBColor objects)
    C_DARK_TEAL = RGBColor(0x00, 0x38, 0x41)   # Titles, icons
    C_TEAL = RGBColor(0x00, 0x7B, 0x85)         # Brand accent, section numbers
    C_BRIGHT_CYAN = RGBColor(0x00, 0xDF, 0xF8)  # Bullets ONLY
    C_DARK_TEXT = RGBColor(0x20, 0x20, 0x20)     # Body text
    C_MEDIUM_GRAY = RGBColor(0x93, 0x97, 0x93)  # Page numbers, dates
    C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)         # Backgrounds

    # Slide dimensions
    SLIDE_WIDTH = Inches(13.33)
    SLIDE_HEIGHT = Inches(7.5)

# Logo positions (inches) — per brand system dimensions.md
LOGO_SMALL = {'x': 0.374, 'y': 7.071, 'w': 0.822, 'h': 0.236}
LOGO_LARGE = {'x': 0.374, 'y': 6.271, 'w': 2.191, 'h': 0.630}
LOGO_BACK = {'x': 5.44, 'y': 2.98, 'w': 2.45, 'h': 1.54}

# Page number position (equal margins from right and bottom edges)
PAGE_NUM = {'x': 12.71, 'y': 7.1554, 'w': 0.33, 'h': 0.152}


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------

def _add_textbox(slide, x, y, w, h, text, *,
                 font_size=14, color=None, bold=False,
                 align=None, font_name=FONT):
    """Add a text box with NBG defaults: margin=0, valign=top, Aptos."""
    if color is None:
        color = C_DARK_TEXT
    if align is None:
        align = PP_ALIGN.LEFT
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    tf.margin_right = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.auto_size = None
    # python-pptx uses MSO_ANCHOR for vertical alignment
    tf.paragraphs[0].alignment = align

    # Set vertical anchor via XML (python-pptx MSO_ANCHOR)
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('anchor', 't')  # top

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name

    return txBox


def _add_bumper_pill(slide, text, x=0.37, y=0.35):
    """Add bumper as a filled rounded-rect pill with white text (NBG pattern).

    Pill: 1.3x0.3, fill #007B85, text 9pt Bold white ALL CAPS.
    """
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(1.3), Inches(0.3),
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = C_TEAL
    pill.line.fill.background()
    pill.shadow.inherit = False

    # Reduce corner rounding
    pill.adjustments[0] = 0.25

    tf = pill.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_top = 0
    tf.margin_right = 0
    tf.margin_bottom = 0
    tf.word_wrap = False

    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.size = Pt(9)
    p.font.color.rgb = C_WHITE
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.LEFT

    # Vertical center
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('anchor', 'ctr')

    return pill


def _add_logo(slide, logo_type='small'):
    """Add NBG logo to slide.

    logo_type: 'small' (content), 'large' (cover/divider), 'back_cover'
    """
    if logo_type == 'back_cover':
        pos = LOGO_BACK
        path = BACK_COVER_LOGO_PNG
    elif logo_type == 'large':
        pos = LOGO_LARGE
        path = LOGO_PNG
    else:
        pos = LOGO_SMALL
        path = LOGO_PNG

    if not path.exists():
        return  # Skip if logo file missing

    slide.shapes.add_picture(
        str(path),
        Inches(pos['x']), Inches(pos['y']),
        Inches(pos['w']), Inches(pos['h']),
    )


def _add_page_number(slide, number):
    """Add page number to bottom-right corner."""
    _add_textbox(
        slide,
        PAGE_NUM['x'], PAGE_NUM['y'], PAGE_NUM['w'], PAGE_NUM['h'],
        str(number),
        font_size=10, color=C_MEDIUM_GRAY, align=PP_ALIGN.RIGHT,
    )


def _add_bullets(slide, x, y, w, h, points, *, font_size=14):
    """Add bullet points with NBG styling: cyan bullet char, Aptos text."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    tf.margin_right = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.auto_size = None

    # Vertical anchor: top
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('anchor', 't')

    for i, point_text in enumerate(points):
        if isinstance(point_text, dict):
            point_text = point_text.get('text', str(point_text))

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = str(point_text)
        p.font.size = Pt(font_size)
        p.font.color.rgb = C_DARK_TEXT
        p.font.name = FONT
        p.alignment = PP_ALIGN.LEFT

        # NBG bullet styling via OOXML: • in Arial, Bright Cyan #00DFF8
        pPr = p._p.get_or_add_pPr()
        pPr.set('lvl', '0')

        buFont = SubElement(pPr, qn('a:buFont'))
        buFont.set('typeface', FONT_BULLET)

        buClr = SubElement(pPr, qn('a:buClr'))
        srgbClr = SubElement(buClr, qn('a:srgbClr'))
        srgbClr.set('val', '00DFF8')

        buChar = SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '\u2022')

        # Paragraph spacing: 14pt before (except first)
        if i > 0:
            spcBef = SubElement(pPr, qn('a:spcBef'))
            spcPts = SubElement(spcBef, qn('a:spcPts'))
            spcPts.set('val', '1400')

    return txBox


# NBG chart color palette (hex strings for series coloring, per charts.md)
CHART_COLORS = ['00ADBF', '003841', '007B85', '939793', 'BEC1BE', '00DFF8']

# Light gray for axis lines and borders
C_LIGHT_GRAY_HEX = 'BEC1BE'


def _style_chart(chart):
    """Apply NBG brand styling to a chart: no title, hidden value axis,
    no gridlines, no plot border, data labels on bars, Aptos throughout."""
    # No chart title — use slide title textbox instead
    chart.has_title = False

    # No legend by default (single-series); caller enables if multi-series
    chart.has_legend = False

    # Category axis: visible, Aptos 12pt, light gray line (per charts.md)
    cat_ax = chart.category_axis
    cat_ax.tick_labels.font.size = Pt(12)
    cat_ax.tick_labels.font.name = FONT
    cat_ax.tick_labels.font.color.rgb = C_DARK_TEXT
    cat_ax.has_major_gridlines = False
    cat_ax.has_minor_gridlines = False
    cat_ax.format.line.fill.solid()
    cat_ax.format.line.fill.fore_color.rgb = RGBColor.from_string(C_LIGHT_GRAY_HEX)

    # Value axis: HIDDEN (per brand spec: valAxisHidden: true)
    val_ax = chart.value_axis
    val_ax.visible = False
    val_ax.has_major_gridlines = False
    val_ax.has_minor_gridlines = False

    # Plot area: no border, no fill
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(11)
    data_labels.font.name = FONT
    data_labels.font.bold = True
    data_labels.font.color.rgb = C_DARK_TEXT
    data_labels.show_value = True
    data_labels.show_category_name = False
    data_labels.show_series_name = False

    # Color each series with NBG palette
    for i, series in enumerate(plot.series):
        color_hex = CHART_COLORS[i % len(CHART_COLORS)]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = RGBColor.from_string(color_hex)


def create_chart_slide(prs, content, page_number, chart_type='bar'):
    """Create a chart slide with NBG styling.

    Bumper:    Pill at (0.37, 0.35)
    Title:     24pt Dark Teal at (0.37, 0.75)
    Chart:     at (0.37, 1.3) — 12.59" x 5.0"
    Logo + page number.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_WHITE

    bumper = content.get('bumper') or content.get('hyper_title')
    if bumper:
        _add_bumper_pill(slide, bumper)

    title_y = 0.75 if bumper else 0.5
    if content.get('title'):
        _add_textbox(slide, 0.37, title_y, 12.59, 0.4, content['title'],
                     font_size=24, color=C_DARK_TEAL, bold=False)

    # Build chart data
    chart_data = CategoryChartData()
    categories = content.get('categories', [])
    series_list = content.get('series', [])
    chart_data.categories = categories

    for s in series_list:
        chart_data.add_series(s.get('name', ''), s.get('values', []))

    # Map chart_type string to XL_CHART_TYPE
    type_map = {
        'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'bar_stacked': XL_CHART_TYPE.COLUMN_STACKED,
        'bar_horizontal': XL_CHART_TYPE.BAR_CLUSTERED,
        'line': XL_CHART_TYPE.LINE_MARKERS,
        'pie': XL_CHART_TYPE.PIE,
    }
    xl_type = type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    chart_y = 1.3 if bumper else 1.1
    chart_h = 5.0 if bumper else 5.2
    chart_frame = slide.shapes.add_chart(
        xl_type, Inches(0.37), Inches(chart_y), Inches(12.59), Inches(chart_h),
        chart_data,
    )

    chart = chart_frame.chart
    _style_chart(chart)

    # Bar gap (per brand: barGapWidthPct: 35)
    chart.plots[0].gap_width = 35

    # Enable legend only for multi-series charts
    if len(series_list) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
        chart.legend.font.name = FONT
        chart.legend.font.color.rgb = C_DARK_TEXT

    _add_logo(slide, 'small')
    _add_page_number(slide, page_number)
    return slide


def create_waterfall_slide(prs, content, page_number):
    """Create a waterfall chart slide using stacked bars.

    python-pptx has no native waterfall type, so we simulate it
    with an invisible base series + positive/negative series.

    Bumper:    Pill at (0.37, 0.35)
    Title:     24pt Dark Teal at (0.37, 0.75)
    Chart:     at (0.37, 1.3) — 12.59" x 5.0"
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_WHITE

    bumper = content.get('bumper') or content.get('hyper_title')
    if bumper:
        _add_bumper_pill(slide, bumper)

    title_y = 0.75 if bumper else 0.5
    if content.get('title'):
        _add_textbox(slide, 0.37, title_y, 12.59, 0.4, content['title'],
                     font_size=24, color=C_DARK_TEAL, bold=False)

    # Waterfall data: list of {label, value} items
    # First and last are totals, middle items are deltas
    items = content.get('waterfall_items', [])
    if not items:
        _add_logo(slide, 'small')
        _add_page_number(slide, page_number)
        return slide

    categories = [item.get('label', '') for item in items]
    values = [item.get('value', 0) for item in items]

    # Build invisible base + increase + decrease series
    base = []
    increase = []
    decrease = []
    running = 0

    for i, val in enumerate(values):
        is_total = items[i].get('total', False)
        if i == 0 or is_total:
            # Total bars start from 0
            base.append(0)
            increase.append(val)
            decrease.append(0)
            running = val
        elif val >= 0:
            base.append(running)
            increase.append(val)
            decrease.append(0)
            running += val
        else:
            base.append(running + val)
            increase.append(0)
            decrease.append(abs(val))
            running += val

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('Base', base)
    chart_data.add_series('Increase', increase)
    chart_data.add_series('Decrease', decrease)

    chart_y = 1.3 if bumper else 1.1
    chart_h = 5.0 if bumper else 5.2
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        Inches(0.37), Inches(chart_y), Inches(12.59), Inches(chart_h),
        chart_data,
    )

    chart = chart_frame.chart
    chart.has_title = False
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 100
    plot.overlap = 100

    # Base series: invisible (no fill, no border, no data labels)
    base_series = plot.series[0]
    base_series.format.fill.background()
    base_series.format.line.fill.background()

    # Increase series: NBG Cyan (per charts.md waterfall spec)
    inc_series = plot.series[1]
    inc_series.format.fill.solid()
    inc_series.format.fill.fore_color.rgb = RGBColor.from_string('00ADBF')
    inc_series.format.line.fill.background()

    # Decrease series: NBG Red (per charts.md/ooxml-charts.md)
    dec_series = plot.series[2]
    dec_series.format.fill.solid()
    dec_series.format.fill.fore_color.rgb = RGBColor.from_string('AA0028')
    dec_series.format.line.fill.background()

    # Data labels: show values on Increase and Decrease series
    for s in [inc_series, dec_series]:
        s.has_data_labels = True
        s.data_labels.font.size = Pt(10)
        s.data_labels.font.name = FONT
        s.data_labels.font.bold = True
        s.data_labels.font.color.rgb = C_DARK_TEXT
        s.data_labels.show_value = True
        s.data_labels.show_category_name = False
        s.data_labels.show_series_name = False

    # Category axis: visible, Aptos, light gray line, no tick marks
    cat_ax = chart.category_axis
    cat_ax.tick_labels.font.size = Pt(9)
    cat_ax.tick_labels.font.name = FONT
    cat_ax.tick_labels.font.color.rgb = C_DARK_TEXT
    cat_ax.has_major_gridlines = False
    cat_ax.has_minor_gridlines = False
    cat_ax.format.line.fill.solid()
    cat_ax.format.line.fill.fore_color.rgb = RGBColor.from_string(C_LIGHT_GRAY_HEX)

    # Value axis: HIDDEN
    val_ax = chart.value_axis
    val_ax.visible = False
    val_ax.has_major_gridlines = False
    val_ax.has_minor_gridlines = False

    _add_logo(slide, 'small')
    _add_page_number(slide, page_number)
    return slide


# ---------------------------------------------------------------------------
# Slide creation functions
# ---------------------------------------------------------------------------

def create_cover_slide(prs, content):
    """Create cover slide with NBG typography hierarchy.

    Title:    48pt Dark Teal at (0.37, 1.39)
    Subtitle: 36pt NBG Teal at (0.37, 2.27)
    Location: 14pt Dark Teal at (0.37, 4.58)
    Date:     14pt Medium Gray at (0.37, 4.97)
    Logo:     Large logo at bottom-left (covers use large)
    NO page number.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_WHITE

    if content.get('title'):
        _add_textbox(slide, 0.37, 1.39, 7.86, 1.00,
                     content['title'],
                     font_size=48, color=C_DARK_TEAL)

    if content.get('subtitle'):
        _add_textbox(slide, 0.37, 2.27, 7.86, 0.80,
                     content['subtitle'],
                     font_size=36, color=C_TEAL)

    if content.get('location'):
        _add_textbox(slide, 0.37, 4.58, 4, 0.4,
                     content['location'],
                     font_size=14, color=C_DARK_TEAL)

    if content.get('date'):
        _add_textbox(slide, 0.37, 4.97, 4, 0.4,
                     content['date'],
                     font_size=14, color=C_MEDIUM_GRAY)

    _add_logo(slide, 'large')
    return slide


def create_divider_slide(prs, content):
    """Create section divider slide.

    Number: 60pt NBG Teal at (0.37, 2.84)
    Title:  48pt Dark Teal at (1.86, 2.84)
    Logo:   Large logo at bottom-left (dividers use large)
    NO page number.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_WHITE

    number = content.get('number', '')
    if isinstance(number, int):
        number = str(number).zfill(2)

    if number:
        _add_textbox(slide, 0.37, 2.84, 1.2, 1.0,
                     number,
                     font_size=60, color=C_TEAL)

    if content.get('title'):
        _add_textbox(slide, 1.86, 2.84, 9.5, 1.0,
                     content['title'],
                     font_size=48, color=C_DARK_TEAL)

    _add_logo(slide, 'large')
    return slide


def create_content_slide(prs, content, page_number):
    """Create content slide with optional bumper pill and bullet points.

    Bumper:  Pill shape at (0.37, 0.35) — 1.3x0.3, #007B85, 9pt Bold white
    Title:   24pt Dark Teal Regular at (0.37, 0.75)
    Bullets: 14pt Dark Text at (0.37, 1.3), cyan bullets
    Logo:    Small logo at bottom-left
    Page number at bottom-right.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_WHITE

    # Bumper pill (optional section tag above title)
    bumper = content.get('bumper') or content.get('hyper_title')
    if bumper:
        _add_bumper_pill(slide, bumper)

    # Action title — 24pt Regular (NOT bold, NOT SemiBold)
    title_y = 0.75 if bumper else 0.5
    if content.get('title'):
        _add_textbox(slide, 0.37, title_y, 12.59, 0.4,
                     content['title'],
                     font_size=24, color=C_DARK_TEAL, bold=False)

    # Body content — bullet points
    points = content.get('points') or content.get('paragraphs') or []
    body_y = 1.3 if bumper else 1.1
    if points:
        _add_bullets(slide, 0.37, body_y, 12.59, 5.0, points)

    _add_logo(slide, 'small')
    _add_page_number(slide, page_number)
    return slide


def create_contents_slide(prs, sections, page_number):
    """Create table of contents slide.

    Header: 32pt Dark Teal Bold at (0.37, 0.36)
    Items:  Number (18pt Teal) + Title (16pt Dark Teal Bold) + Desc (12pt)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_WHITE

    _add_textbox(slide, 0.37, 0.36, 10, 0.70,
                 'Contents',
                 font_size=32, color=C_DARK_TEAL, bold=True)

    for i, section in enumerate(sections):
        y = 1.48 + (i * 0.85)

        number = section.get('number', str(i + 1).zfill(2))
        _add_textbox(slide, 0.37, y, 0.60, 0.60,
                     str(number),
                     font_size=18, color=C_TEAL, bold=True)

        if section.get('title'):
            _add_textbox(slide, 1.10, y, 8, 0.35,
                         section['title'],
                         font_size=16, color=C_DARK_TEAL, bold=True)

        if section.get('description'):
            _add_textbox(slide, 1.10, y + 0.35, 8, 0.30,
                         section['description'],
                         font_size=12, color=RGBColor(0x59, 0x59, 0x59))

    _add_logo(slide, 'small')
    _add_page_number(slide, page_number)
    return slide


def create_back_cover_slide(prs):
    """Create back cover slide.

    IMPORTANT:
    - NO "Thank You" text
    - NO corner logo
    - NO page number
    - White background with centered oval NBG building logo ONLY
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_WHITE

    _add_logo(slide, 'back_cover')
    return slide


# ---------------------------------------------------------------------------
# Catalog and type normalization (kept for backward compatibility + tests)
# ---------------------------------------------------------------------------

def load_catalog():
    """Load the slide catalog."""
    with open(CATALOG_PATH, 'r') as f:
        return yaml.safe_load(f)


def resolve_slide_index(catalog, template, slide_type):
    """Resolve a slide type path to template index (for legacy compatibility)."""
    parts = slide_type.split('/')
    slides = catalog['templates'][template]['slides']

    current = slides
    for part in parts:
        if isinstance(current, dict) and part in current:
            current = current[part]
        else:
            raise ValueError(f"Unknown slide type: {slide_type}")

    if isinstance(current, int):
        return current
    raise ValueError(f"Slide type {slide_type} does not resolve to an index")


def normalize_slide_type(slide_type, recommended_visual=None):
    """Convert McKinsey slide types to catalog paths.

    Handles both:
    - Simple: "covers/simple_white" (passes through)
    - McKinsey: "cover", "content", "divider", "back_cover"

    Uses recommended_visual hint when available for smarter selection.
    """
    if '/' in slide_type:
        return slide_type

    visual_mapping = {
        'bar_chart': 'charts/bar_single',
        'pie_chart': 'charts/pie_single',
        'line_chart': 'charts/line_single',
        'waterfall_chart': 'charts/bar_single',
        'comparison_chart': 'charts/bar_dual',
        'kpi_dashboard': 'content/text_only',
        'numbered_infographic': 'content/text_only',
        'timeline': 'content/text_only',
        'process': 'content/text_only',
        'funnel': 'content/text_only',
        'table': 'tables/half_page',
        'comparison_table': 'tables/comparison',
        'icons': 'content/text_only',
        'none': 'content/text_only',
    }

    if recommended_visual and recommended_visual in visual_mapping:
        return visual_mapping[recommended_visual]

    type_mapping = {
        'cover': 'covers/quiet',
        'divider': 'dividers/quiet_white',
        'content': 'content/text_only',
        'chart': 'charts/bar_dual',
        'infographic': 'infographics/numbered_6',
        'table': 'tables/half_page',
        'thankyou': 'back_covers/quiet',
        'back_cover': 'back_covers/quiet',
        'summary': 'content/text_with_bullets',
    }

    return type_mapping.get(slide_type, 'content/text_only')


# ---------------------------------------------------------------------------
# Slide type detection
# ---------------------------------------------------------------------------

def _classify_slide(slide_def):
    """Classify a slide definition into a build category.

    Returns one of: 'cover', 'divider', 'contents', 'content', 'back_cover'
    """
    raw_type = slide_def.get('type', 'content')

    # Direct match on raw type
    if raw_type in ('cover', 'covers') or raw_type.startswith('covers/'):
        return 'cover'
    if raw_type in ('back_cover', 'thankyou') or raw_type.startswith('back_covers/'):
        return 'back_cover'
    if raw_type in ('divider',) or raw_type.startswith('dividers/'):
        return 'divider'
    if raw_type in ('contents', 'toc'):
        return 'contents'
    if raw_type in ('chart', 'bar_chart', 'line_chart', 'pie_chart'):
        return 'chart'
    if raw_type in ('waterfall', 'waterfall_chart'):
        return 'waterfall'

    # Check recommended_visual for chart hints
    visual = slide_def.get('recommended_visual', '')
    if visual in ('bar_chart', 'comparison_chart'):
        return 'chart'
    if visual == 'waterfall_chart':
        return 'waterfall'

    # Everything else is content (text, infographics, tables)
    return 'content'


# ---------------------------------------------------------------------------
# Main build function
# ---------------------------------------------------------------------------

def build_presentation(outline_path, output_path):
    """Build NBG presentation from YAML outline.

    Creates every slide from scratch with python-pptx. No template extraction.
    White backgrounds, exact NBG positioning, proper typography.
    """
    outline_path = Path(outline_path).expanduser().resolve()
    output_path = Path(output_path).expanduser().resolve()

    with open(outline_path, 'r') as f:
        outline = yaml.safe_load(f)

    # Support both McKinsey and simple formats
    if 'presentation' in outline:
        presentation = outline['presentation']
        slides = presentation.get('slides', [])
        template = outline.get('template', 'GR')
        print(f"\n[McKinsey Quality] Using Pyramid Principle structure")
        if presentation.get('main_recommendation'):
            print(f"  Main recommendation: {presentation['main_recommendation'][:60]}...")
    else:
        template = outline.get('template', 'GR')
        slides = outline.get('slides', [])

    if not slides:
        raise ValueError("No slides defined in outline")

    print(f"\nNBG Presentation Builder")
    print(f"{'='*50}")
    print(f"Template: {template}")
    print(f"Slides: {len(slides)}")
    print(f"Output: {output_path}")
    print(f"{'='*50}\n")

    # Create presentation from scratch
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    page_number = 0

    for slide_def in slides:
        category = _classify_slide(slide_def)
        content = slide_def.get('content', {})
        raw_type = slide_def.get('type', 'content')

        if category == 'cover':
            create_cover_slide(prs, content)
            print(f"  [cover] {content.get('title', '')[:50]}")

        elif category == 'divider':
            create_divider_slide(prs, content)
            print(f"  [divider] {content.get('number', '')} {content.get('title', '')[:50]}")

        elif category == 'contents':
            page_number += 1
            sections = content.get('sections', [])
            create_contents_slide(prs, sections, page_number)
            print(f"  [contents] {len(sections)} sections")

        elif category == 'chart':
            page_number += 1
            chart_subtype = raw_type.replace('_chart', '') if '_chart' in raw_type else 'bar'
            create_chart_slide(prs, content, page_number, chart_type=chart_subtype)
            print(f"  [chart:{chart_subtype}] {content.get('title', '')[:50]}")

        elif category == 'waterfall':
            page_number += 1
            create_waterfall_slide(prs, content, page_number)
            print(f"  [waterfall] {content.get('title', '')[:50]}")

        elif category == 'back_cover':
            create_back_cover_slide(prs)
            print(f"  [back_cover]")

        else:  # content
            page_number += 1
            create_content_slide(prs, content, page_number)
            visual = slide_def.get('recommended_visual', '')
            label = f" ({visual})" if visual else ''
            print(f"  [content{label}] {content.get('title', '')[:50]}")

    prs.save(str(output_path))

    print(f"\n{'='*50}")
    print(f"Created: {output_path}")
    print(f"{'='*50}\n")

    # Validate
    validate_script = SCRIPT_DIR / 'nbg_validate.py'
    if validate_script.exists():
        print("Validating...")
        subprocess.run([sys.executable, str(validate_script), str(output_path)])


def main():
    if len(sys.argv) < 3:
        print("NBG Presentation Builder")
        print("=" * 40)
        print("\nUsage: python nbg_build.py <outline.yaml> <output.pptx>")
        print("\nCreates presentations from scratch with NBG brand guidelines.")
        print("White backgrounds, Aptos font, pixel-perfect positioning.")
        print("\nExample outline.yaml:")
        print("""
template: GR

presentation:
  title: "Presentation Title"
  audience: "Board of Directors"
  main_recommendation: "Lead with the answer"

  slides:
    - type: cover
      content:
        title: "Presentation Title"
        subtitle: "Subtitle"
        date: "March 2026"

    - type: content
      key_message: "The key insight"
      so_what: "Why it matters"
      content:
        title: "Action title that tells the story"
        bumper: "SECTION TAG"
        points:
          - "First bullet point"
          - "Second bullet point"

    - type: back_cover
      content: {}
""")
        sys.exit(1)

    try:
        build_presentation(sys.argv[1], sys.argv[2])
    except Exception as e:
        print(f"\nError: {e}")
        sys.exit(1)


if __name__ == '__main__':
    main()
